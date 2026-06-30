"""
Create PowerPoint Presentation from SPEC.md content
Professional Academic Style for AeroCast VAA System
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color scheme - Professional Academic
NAVY = RGBColor(0x1F, 0x4E, 0x79)
LIGHT_BLUE = RGBColor(0x5B, 0x9B, 0xD5)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

# Paths
DIAGRAM_DIR = os.path.dirname(__file__)

def set_shape_fill(shape, color):
    """Set solid fill color for a shape"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def add_title_box(slide, text, left, top, width, height):
    """Add a title text box"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_content_box(slide, text_lines, left, top, width, height):
    """Add a content text box with bullet points"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True

    for i, line in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
        if line.startswith("• "):
            p.text = line[2:]
            p.level = 0
        elif line.startswith("  - "):
            p.text = line[4:]
            p.level = 1
    return shape

def add_image_to_slide(slide, img_path, left, top, width, height):
    """Add an image to slide"""
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, left, top, width, height)
        return True
    else:
        print(f"Warning: Image not found: {img_path}")
        return False

def add_slide_number(slide, number, total):
    """Add slide number at bottom"""
    left = Inches(0.5)
    top = Inches(7.0)
    width = Inches(9)
    height = Inches(0.3)

    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = f"{number} / {total}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
    p.alignment = PP_ALIGN.CENTER

def create_title_slide(prs, slide_num, total):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    set_shape_fill(shape, NAVY)
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AEROCAST VAA SYSTEM"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Mô hình dự báo NeuralProphet cho Air Traffic"
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Author info
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    tf = author_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Sinh viên: Lương Minh Khôi"
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "GVHD: Th.S [Tên thầy]"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Ngày: 29/05/2026"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p.alignment = PP_ALIGN.CENTER

def create_toc_slide(prs, slide_num, total):
    """Slide 2: Table of Contents"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_box(slide, "MỤC LỤC", Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))

    toc_items = [
        "1. Tổng quan hệ thống AeroCast",
        "2. Bài toán dự báo là gì?",
        "3. Tại sao không dùng công thức toán đơn giản?",
        "4. Neural Network hoạt động thế nào?",
        "5. NeuralProphet - Kiến trúc và hoạt động",
        "6. Giải thích Hyperparameters",
        "7. So sánh các phương pháp dự báo",
        "8. Tại sao chọn NeuralProphet?",
        "9. Pipeline hoàn chỉnh",
        "10. Kết luận & Q&A"
    ]

    toc_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(8), Inches(5.5))
    tf = toc_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(toc_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)

    add_slide_number(slide, slide_num, total)

def create_diagram_slide(prs, slide_num, total, title, img_name, description=None):
    """Generic diagram slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    add_title_box(slide, title, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))

    # Diagram
    img_path = os.path.join(DIAGRAM_DIR, img_name)
    add_image_to_slide(slide, img_path, Inches(0.5), Inches(1.3), Inches(9), Inches(5.2))

    # Description if provided
    if description:
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = description
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p.alignment = PP_ALIGN.CENTER

    add_slide_number(slide, slide_num, total)

def create_content_slide(prs, slide_num, total, title, bullet_points, two_columns=False):
    """Generic content slide with bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_box(slide, title, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))

    if two_columns and len(bullet_points) > 6:
        # Split into two columns
        half = len(bullet_points) // 2
        left_points = bullet_points[:half]
        right_points = bullet_points[half:]

        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.3), Inches(5.5))
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, point in enumerate(left_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + point if not point.startswith("•") else point
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(8)

        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.5))
        tf = right_box.text_frame
        tf.word_wrap = True
        for i, point in enumerate(right_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + point if not point.startswith("•") else point
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(8)
    else:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + point if not point.startswith("•") else point
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(10)

    add_slide_number(slide, slide_num, total)

def create_comparison_slide(prs, slide_num, total, title):
    """Slide with comparison table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_box(slide, title, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))

    # Add comparison image
    img_path = os.path.join(DIAGRAM_DIR, "diagram_07_comparison.png")
    add_image_to_slide(slide, img_path, Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))

    add_slide_number(slide, slide_num, total)

def create_code_slide(prs, slide_num, total, title, code_lines):
    """Slide with code block"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_box(slide, title, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))

    # Code background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(9), Inches(5.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
    shape.line.fill.background()

    # Code text
    code_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
    tf = code_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(code_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.name = "Courier New"
        p.font.color.rgb = RGBColor(0xA9, 0xB7, 0xC6)

    add_slide_number(slide, slide_num, total)

def create_conclusion_slide(prs, slide_num, total):
    """Slide 20: Conclusion"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_box(slide, "KẾT LUẬN", Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))

    conclusions = [
        "NeuralProphet là lựa chọn TỐI ƯU cho bài toán dự báo air traffic",
        "Phù hợp với dataset nhỏ (~20-36 điểm dữ liệu)",
        "Hỗ trợ native: exogenous variables (weather, buffer)",
        "Uncertainty estimation với quantiles [0.1, 0.9]",
        "Cấu hình hiện tại GẦN NHƯ TỐI ƯU cho dữ liệu hiện có",
        "Có thể cải thiện khi thu thập thêm dữ liệu"
    ]

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(conclusions):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "✓ " + point
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)  # Green
        p.space_after = Pt(14)

    add_slide_number(slide, slide_num, total)

def create_qa_slide(prs, slide_num, total):
    """Slide 21: Q&A"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    set_shape_fill(shape, NAVY)
    shape.line.fill.background()

    # Q&A Text
    qa_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    tf = qa_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Q & A"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Cảm ơn thầy đã lắng nghe!"
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    add_slide_number(slide, slide_num, total)

def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    total_slides = 21
    slide_num = 1

    # Slide 1: Title
    create_title_slide(prs, slide_num, total_slides)
    slide_num += 1

    # Slide 2: TOC
    create_toc_slide(prs, slide_num, total_slides)
    slide_num += 1

    # Slide 3-4: System Overview
    create_diagram_slide(prs, slide_num, total_slides,
                        "1. Tổng quan hệ thống AeroCast",
                        "diagram_01_system_overview.png",
                        "Hình 1: Luồng dữ liệu từ Radar → Database → AI Model → Prediction")
    slide_num += 1

    create_diagram_slide(prs, slide_num, total_slides,
                        "1. Đầu vào và đầu ra của mô hình",
                        "diagram_02_io.png",
                        "Hình 2: Input (60 phút quá khứ) → Output (30 phút dự đoán)")
    slide_num += 1

    # Slide 5: Bài toán dự báo
    create_content_slide(prs, slide_num, total_slides, "2. Bài toán dự báo là gì?", [
        "Dự báo = Dự đoán tương lai dựa trên quá khứ",
        "Ví dụ đời thường: Dự báo thời tiết, kẹt xe, lượng khách",
        "Trong Air Traffic: Dự đoán số máy bay trong sector trong 30 phút tới",
        "Mục tiêu: Giúp controller chuẩn bị và điều phối tốt hơn"
    ])
    slide_num += 1

    # Slide 6: Tại sao dự báo khó
    create_content_slide(prs, slide_num, total_slides, "2. Tại sao dự báo khó?", [
        "Quá khứ không hoàn toàn lặp lại - mỗi ngày khác nhau",
        "Có nhiều yếu tố ảnh hưởng: thời tiết, sự kiện, mùa, giờ cao điểm",
        "Sự bất định (Uncertainty) - không thể đoán chính xác 100%",
        "Cần ước lượng khoảng dao động thay vì 1 con số cố định"
    ])
    slide_num += 1

    # Slide 7-8: Neural Network
    create_diagram_slide(prs, slide_num, total_slides,
                        "4. Neural Network hoạt động thế nào?",
                        "diagram_03_neural_network.png",
                        "Hình 3: Kiến trúc Neural Network với 2 hidden layers [32, 32]")
    slide_num += 1

    create_content_slide(prs, slide_num, total_slides, "4. Neural Network - Cách 'học'", [
        "Bước 1: Khởi tạo weights ngẫu nhiên",
        "Bước 2: Tính sai số (Error) - so sánh dự đoán với thực tế",
        "Bước 3: Điều chỉnh weights để sai số giảm (Backpropagation)",
        "Bước 4: Lặp lại nhiều epoch cho đến khi sai số đủ nhỏ",
        "'Học' = Điều chỉnh trọng số để dự đoán chính xác hơn"
    ])
    slide_num += 1

    # Slide 9-12: NeuralProphet Architecture
    create_diagram_slide(prs, slide_num, total_slides,
                        "5. NeuralProphet - Tổng quan kiến trúc",
                        "diagram_04_ar_net.png",
                        "Hình 4: AR-Net là điểm khác biệt chính với Facebook Prophet")
    slide_num += 1

    create_content_slide(prs, slide_num, total_slides, "5. NeuralProphet - Các thành phần", [
        "TREND: Xu hướng dài hạn (tăng/giảm theo thời gian)",
        "SEASONALITY: Pattern lặp lại (theo ngày, tuần)",
        "AR-NET: Autoregressive Neural Network - dùng quá khứ dự đoán tương lai",
        "REGRESSORS: Biến ngoại sinh (weather, buffer counts)",
        "QUANTILES: Ước lượng uncertainty (khoảng dao động)"
    ])
    slide_num += 1

    create_diagram_slide(prs, slide_num, total_slides,
                        "5. Seasonality - Pattern theo chu kỳ",
                        "diagram_08_seasonality.png",
                        "Hình 5: Daily (theo giờ) và Weekly (theo ngày) seasonality")
    slide_num += 1

    create_diagram_slide(prs, slide_num, total_slides,
                        "5. Quantiles - Ước lượng bất định",
                        "diagram_05_quantiles.png",
                        "Hình 6: Khoảng dao động 80% confidence [y_10, y_90]")
    slide_num += 1

    # Slide 13-14: Hyperparameters
    create_diagram_slide(prs, slide_num, total_slides,
                        "6. Giải thích Hyperparameters",
                        "diagram_09_hyperparams.png",
                        "Hình 7: 4 hyperparameters chính của NeuralProphet")
    slide_num += 1

    create_content_slide(prs, slide_num, total_slides, "6. Hyperparameters - Chi tiết", [
        "n_lags = 12: Dùng 12 điểm dữ liệu quá khứ (60 phút)",
        "n_forecasts = 6: Dự đoán 6 bước tương lai (30 phút)",
        "ar_layers = [32, 32]: Neural network 2 tầng, mỗi tầng 32 neurons",
        "quantiles = [0.1, 0.9]: Khoảng dao động 80% (10th - 90th percentile)"
    ])
    slide_num += 1

    # Slide 15: Comparison
    create_comparison_slide(prs, slide_num, total_slides, "7. So sánh các phương pháp")
    slide_num += 1

    # Slide 16-17: Why NeuralProphet
    create_content_slide(prs, slide_num, total_slides, "8. Tại sao chọn NeuralProphet?", [
        "THIẾT KẾ CHO DATA NHỎ: ~20 điểm là đủ để train",
        "AR-NET: Kết hợp autoregressive + neural network - tốt nhất cả 2 thế giới",
        "Native exogenous support: Thời tiết, buffer counts dễ dàng",
        "Uncertainty estimation: Quantile regression tích hợp sẵn",
        "Interpretable: Component decomposition dễ giải thích với thầy"
    ])
    slide_num += 1

    create_content_slide(prs, slide_num, total_slides, "8. Tại sao KHÔNG dùng phương pháp khác?", [
        "ARIMA: Chỉ học được quan hệ TUYẾN TÍNH, khó dùng với nhiều biến",
        "LSTM: Cần ~1000+ điểm data, dataset nhỏ sẽ OVERFIT nặng",
        "Transformer: Cần ~5000+ điểm data, quá phức tạp cho bài toán này",
        "Prophet (Facebook): Không có AR-Net, không tận dụng được autoregressive"
    ])
    slide_num += 1

    # Slide 18-19: Pipeline
    create_diagram_slide(prs, slide_num, total_slides,
                        "9. Pipeline hoàn chỉnh",
                        "diagram_06_pipeline.png",
                        "Hình 8: 9 bước từ Data Collection → Prediction")
    slide_num += 1

    create_code_slide(prs, slide_num, total_slides, "9. Cấu hình NeuralProphet (Code)", [
        "m = NeuralProphet(",
        "    n_lags=12,",
        "    n_forecasts=6,",
        "    ar_layers=[32, 32],",
        "    quantiles=[0.1, 0.9],",
        "    yearly_seasonality=False,",
        "    weekly_seasonality=True,",
        "    daily_seasonality=True,",
        ")"
    ])
    slide_num += 1

    # Slide 20: Conclusion
    create_conclusion_slide(prs, slide_num, total_slides)
    slide_num += 1

    # Slide 21: Q&A
    create_qa_slide(prs, slide_num, total_slides)

    # Save
    output_path = os.path.join(DIAGRAM_DIR, "AeroCast_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    main()
